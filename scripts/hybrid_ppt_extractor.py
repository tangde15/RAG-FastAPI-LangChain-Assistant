
"""
Hybrid-PPT-Extractor: Windows可用，支持结构化+图片OCR+备注全量提取
依赖：unstructured[local-inference]、python-pptx、pillow、pandas、paddlepaddle、paddleocr
"""
from unstructured.partition.pptx import partition_pptx
from pptx import Presentation
from paddleocr import PaddleOCR
import io
from PIL import Image
import threading
import numpy as np
import inspect
import tempfile
import os
try:
    import cv2
except Exception:
    cv2 = None

# ============ PaddleOCR 单例模式（不会重复初始化） ============
_ocr_instance = None
_ocr_lock = threading.Lock()

def get_ocr():
    global _ocr_instance
    if _ocr_instance is None:
        with _ocr_lock:
            if _ocr_instance is None:   # 双重检查，线程安全
                print("[OCR] 初始化 PaddleOCR（仅首次，增强小字识别）...")
                _ocr_instance = PaddleOCR(
                    use_textline_orientation=True,  # 文本方向校正
                    lang='ch',
                    text_det_thresh=0.2,    # 默认0.3，调低后更敏感
                    text_det_box_thresh=0.4, # 默认0.5，降低文本框过滤阈值
                    text_det_unclip_ratio=2.0, # 扩大文本框范围，避免小字被截断
                    text_rec_input_shape="3, 48, 320", # 调整识别图像尺寸，适配小字
                    text_recognition_batch_size=1  # 单批次识别，提升小字识别准确率
                )
    return _ocr_instance
# ==================================================================

def extract_text_unstructured(filepath):
    return partition_pptx(
        filepath,
        strategy="hi_res",
        infer_table_structure=True,
        extract_image_block=True
    )

def extract_text_pythonpptx(filepath):
    prs = Presentation(filepath)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                texts.append("备注：" + notes)
    return texts

def extract_ocr_images(filepath):
    prs = Presentation(filepath)
    ocr_texts = []
    ocr = get_ocr()   # ⭐ 只使用单例 OCR

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # PICTURE
                img_bytes = shape.image.blob
                img = Image.open(io.BytesIO(img_bytes))

                # 将 PIL Image 转为 numpy.ndarray（BGR），PaddleOCR 期望 numpy 或本地路径
                try:
                    np_img = None
                    if cv2 is not None:
                        try:
                            # 如果是 PIL Image，先转换为 RGB 再转为 BGR numpy
                            if isinstance(img, Image.Image):
                                np_img = cv2.cvtColor(np.array(img), cv2.COLOR_RGB2BGR)
                            else:
                                # 非 PIL，尝试直接转换
                                np_img = np.array(img)
                        except Exception:
                            np_img = None

                    # 决定使用 numpy 图像还是临时文件路径
                    call_arg = None
                    temp_path = None
                    if np_img is not None:
                        call_arg = np_img
                    else:
                        # 回退：保存到临时文件并传路径给 PaddleOCR
                        try:
                            fd, temp_path = tempfile.mkstemp(suffix='.png')
                            os.close(fd)
                            img.save(temp_path)
                            call_arg = temp_path
                        except Exception as e:
                            print(f"[OCR] 无法准备图片用于 OCR: {e}")
                            continue

                    # 调用 OCR：判断方法签名是否支持 cls 参数（只在 ocr.ocr 可用时尝试）
                    result = None
                    try:
                        if hasattr(ocr, 'ocr'):
                            sig = inspect.signature(ocr.ocr)
                            if 'cls' in sig.parameters:
                                result = ocr.ocr(call_arg, cls=True)
                            else:
                                result = ocr.ocr(call_arg)
                        elif hasattr(ocr, 'predict'):
                            sig = inspect.signature(ocr.predict)
                            if 'cls' in sig.parameters:
                                result = ocr.predict(call_arg, cls=True)
                            else:
                                result = ocr.predict(call_arg)
                        else:
                            raise AttributeError('OCR 对象无 ocr 或 predict 方法')
                    finally:
                        # 清理临时文件（如果有）
                        if temp_path and os.path.exists(temp_path):
                            try:
                                os.remove(temp_path)
                            except Exception:
                                pass

                    # 解析结果（兼容所有 PaddleOCR 版本的安全写法）
                    # result 格式通常为: [[[box, (text, score)], ...], ...] 或 [[box, (text, score)], ...]
                    if result:
                        print(f"[OCR] 调试 - result 类型: {type(result)}, 长度: {len(result) if isinstance(result, (list, tuple)) else 'N/A'}")
                        for idx, line in enumerate(result):
                            if not line:
                                continue
                            print(f"[OCR] 调试 - line[{idx}] 类型: {type(line)}, 长度: {len(line) if isinstance(line, (list, tuple)) else 'N/A'}")
                            # line 可能是列表，遍历每个检测项
                            if isinstance(line, list):
                                for item_idx, item in enumerate(line):
                                    try:
                                        print(f"[OCR] 调试 - item[{item_idx}] 类型: {type(item)}, 长度: {len(item) if isinstance(item, (list, tuple)) else 'N/A'}")
                                        # item 格式: [box, (text, score)] 或 [box, text, score] 等
                                        if isinstance(item, (list, tuple)) and len(item) >= 2:
                                            # 标准格式: item[1] 是 (text, score)
                                            text_info = item[1]
                                            if isinstance(text_info, (list, tuple)) and len(text_info) >= 1:
                                                text = text_info[0]
                                                if text and isinstance(text, str):
                                                    ocr_texts.append(text)
                                                    print(f"[OCR] 提取文本: {text}")
                                            elif isinstance(text_info, str):
                                                ocr_texts.append(text_info)
                                                print(f"[OCR] 提取文本: {text_info}")
                                    except Exception as e:
                                        print(f"[OCR] 解析 item 出错: {e}")
                                        continue
                except Exception as e:
                    print(f"[OCR] 识别异常: {e}")
    return ocr_texts

def extract_all(filepath):
    final_text = []

    print("[1] 解析结构化内容 (unstructured)...")
    try:
        elements = extract_text_unstructured(filepath)
        for e in elements:
            txt = str(e).strip()
            if txt:
                final_text.append(txt)
    except Exception as e:
        print("unstructured 解析失败：", e)

    print("[2] 解析遗漏的文本框与备注 (python-pptx)...")
    try:
        extra = extract_text_pythonpptx(filepath)
        final_text.extend(extra)
    except Exception:
        pass

    print("[3] OCR 识别图片文字...")
    try:
        ocr_parts = extract_ocr_images(filepath)
        final_text.extend(ocr_parts)
    except Exception as e:
        print("OCR失败：", e)

    # 去重 & 清洗
    cleaned = []
    seen = set()
    for t in final_text:
        t = t.strip()
        if t and t not in seen:
            seen.add(t)
            cleaned.append(t)

    return "\n".join(cleaned)

if __name__ == "__main__":
    import sys
    if len(sys.argv) < 2:
        print("用法: python hybrid_ppt_extractor.py <your.pptx>")
    else:
        text = extract_all(sys.argv[1])
        print(text)
