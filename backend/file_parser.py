"""
文件解析模块：支持 PDF、DOCX、PPTX、TXT 等格式
"""
import os
from typing import List
from transformers import AutoTokenizer
import re

def read_pdf(path: str) -> str:
    """解析 PDF 文件"""
    try:
        import pdfplumber
        print(f"[文件解析] 开始解析 PDF 文件: {os.path.basename(path)}")
        text = []
        with pdfplumber.open(path) as pdf:
            total_pages = len(pdf.pages)
            print(f"[文件解析] PDF 共 {total_pages} 页")
            for i, page in enumerate(pdf.pages, 1):
                page_text = page.extract_text()
                if page_text:
                    text.append(page_text)
                    print(f"[文件解析] 已解析第 {i}/{total_pages} 页")
        result = "\n".join(text)
        print(f"[文件解析] PDF 解析完成，提取文本 {len(result)} 字符")
        return result
    except ImportError:
        raise ImportError("请安装 pdfplumber: pip install pdfplumber")
    except Exception as e:
        raise Exception(f"PDF 解析失败: {str(e)}")


def read_docx(path: str) -> str:
    """解析 Word 文件"""
    try:
        import docx
        print(f"[文件解析] 开始解析 Word 文件: {os.path.basename(path)}")
        doc = docx.Document(path)
        paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
        result = "\n".join(paragraphs)
        print(f"[文件解析] Word 解析完成，提取 {len(paragraphs)} 个段落，共 {len(result)} 字符")
        return result
    except ImportError:
        raise ImportError("请安装 python-docx: pip install python-docx")
    except Exception as e:
        raise Exception(f"DOCX 解析失败: {str(e)}")


def read_pptx(path: str) -> str:
    """解析 PowerPoint 文件"""
    try:
        from pptx import Presentation
        print(f"[文件解析] 开始解析 PowerPoint 文件: {os.path.basename(path)}")
        prs = Presentation(path)
        texts = []
        total_slides = len(prs.slides)
        print(f"[文件解析] PPT 共 {total_slides} 页")
        for i, slide in enumerate(prs.slides, 1):
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            print(f"[文件解析] 已解析第 {i}/{total_slides} 页")
        result = "\n".join(texts)
        print(f"[文件解析] PPT 解析完成，提取文本 {len(result)} 字符")
        return result
    except ImportError:
        raise ImportError("请安装 python-pptx: pip install python-pptx")
    except Exception as e:
        raise Exception(f"PPTX 解析失败: {str(e)}")


def read_txt(path: str) -> str:
    """解析文本文件"""
    try:
        print(f"[文件解析] 开始解析文本文件: {os.path.basename(path)}")
        with open(path, "r", encoding="utf-8") as f:
            result = f.read()
        print(f"[文件解析] 文本文件解析完成，共 {len(result)} 字符")
        return result
    except UnicodeDecodeError:
        # 尝试其他编码
        try:
            print(f"[文件解析] UTF-8 编码失败，尝试 GBK 编码")
            with open(path, "r", encoding="gbk") as f:
                result = f.read()
            print(f"[文件解析] 文本文件解析完成（GBK），共 {len(result)} 字符")
            return result
        except Exception:
            print(f"[文件解析] GBK 编码失败，尝试 Latin-1 编码")
            with open(path, "r", encoding="latin-1") as f:
                result = f.read()
            print(f"[文件解析] 文本文件解析完成（Latin-1），共 {len(result)} 字符")
            return result
    except Exception as e:
        raise Exception(f"TXT 解析失败: {str(e)}")


def read_file(path: str) -> str:
    """
    根据文件扩展名自动选择解析器
    支持: .pdf, .docx, .pptx, .txt, .md
    """
    ext = os.path.splitext(path)[1].lower()
    
    if ext == '.pdf':
        return read_pdf(path)
    elif ext == '.docx':
        return read_docx(path)
    elif ext == '.pptx':
        return read_pptx(path)
    elif ext in ['.txt', '.md']:
        return read_txt(path)
    else:
        raise ValueError(f"不支持的文件类型: {ext}")


class BGEChunker:
    """BGE 语义切片器：基于 BGE tokenizer 的句子级切片"""
    def __init__(self, 
                 model_name="BAAI/bge-base-zh-v1.5",
                 chunk_size=500,
                 overlap=50):
        print(f"[BGE切片器] 初始化，模型: {model_name}")
        self.tokenizer = AutoTokenizer.from_pretrained(model_name)
        # 强制限制 chunk_size 在 500-1000 范围内
        self.chunk_size = max(500, min(1000, chunk_size))
        self.overlap = overlap
        print(f"[BGE切片器] chunk_size={self.chunk_size} 字符, overlap={overlap} 字符")

    def split_sentences(self, text):
        """中文断句规则"""
        text = re.sub(r"\s+", " ", text)
        sentences = re.split(r"(?<=[。！？，；：\n])", text)
        return [s.strip() for s in sentences if s.strip()]

    def tokenize_len(self, text):
        """计算文本 token 长度"""
        return len(self.tokenizer.encode(text))

    def merge_sentences(self, sentences):
        """合并句子到满足 chunk_size"""
        chunks = []
        current = []

        for sent in sentences:
            current.append(sent)
            token_len = self.tokenize_len("".join(current))

            if token_len > self.chunk_size:
                # 超长 → 分段
                if len(current) > 1:
                    chunks.append("".join(current[:-1]))
                    current = [current[-1]]
                else:
                    # 单句超长，直接加入
                    chunks.append(current[0])
                    current = []

        if current:
            chunks.append("".join(current))

        return chunks

    def add_overlap(self, chunks):
        """重叠处理"""
        final = []
        prev_tokens = []

        for chunk in chunks:
            tokens = self.tokenizer.encode(chunk)
            # 添加重叠
            if prev_tokens and len(prev_tokens) > self.overlap:
                overlap_tokens = prev_tokens[-self.overlap:]
                merged = self.tokenizer.decode(overlap_tokens + tokens, skip_special_tokens=True)
                final.append(merged)
            else:
                final.append(chunk)

            prev_tokens = tokens

        return final

    def chunk(self, text):
        """主入口：文本 → 语义切片"""
        sentences = self.split_sentences(text)
        print(f"[BGE切片器] 分句完成，共 {len(sentences)} 句")
        merged = self.merge_sentences(sentences)
        print(f"[BGE切片器] 合并句子，生成 {len(merged)} 个初步片段")
        final = self.add_overlap(merged)
        print(f"[BGE切片器] 添加重叠，最终 {len(final)} 个片段")
        return final


# 全局切片器实例（避免重复加载 tokenizer）
_chunker = None

def get_chunker():
    """获取全局切片器实例"""
    global _chunker
    if _chunker is None:
        _chunker = BGEChunker(
            model_name="BAAI/bge-base-zh-v1.5",
            chunk_size=500,
            overlap=50
        )
    return _chunker


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[str]:
    """
    将长文本切分成多个片段（基于 BGE 语义切片）
    
    Args:
        text: 原始文本
        chunk_size: 每个片段的最大字符数（默认 500，强制范围 500-1000）
        overlap: 片段之间的重叠字符数（默认 50）
    
    Returns:
        文本片段列表
    """
    if not text or not text.strip():
        return []
    
    # 强制限制 chunk_size 在 500-1000 范围内
    chunk_size = max(500, min(1000, chunk_size))
    
    print(f"[文本切片] 开始语义切片，文本长度: {len(text)} 字符")
    print(f"[文本切片] chunk_size={chunk_size} 字符, overlap={overlap} 字符")
    
    chunker = get_chunker()
    chunks = chunker.chunk(text)
    
    print(f"[文本切片] 语义切片完成，共生成 {len(chunks)} 个片段")
    return chunks


# 支持的文件类型
SUPPORTED_EXTENSIONS = {'.pdf', '.docx', '.pptx', '.txt', '.md'}


def is_supported_file(filename: str) -> bool:
    """检查文件是否支持"""
    ext = os.path.splitext(filename)[1].lower()
    return ext in SUPPORTED_EXTENSIONS
